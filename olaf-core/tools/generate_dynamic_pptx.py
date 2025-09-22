import os
import re
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

def split_bullet_content(content_text, max_bullets=7):
    """Split content if it has more than max_bullets bullet points."""
    lines = content_text.split('\n')
    bullets = [line for line in lines if line.strip().startswith('-') or line.strip().startswith('*')]
    
    if len(bullets) <= max_bullets:
        return [content_text]
    
    # Split bullets into chunks
    bullet_chunks = []
    non_bullet_lines = [line for line in lines if not (line.strip().startswith('-') or line.strip().startswith('*'))]
    
    for i in range(0, len(bullets), max_bullets):
        chunk_bullets = bullets[i:i + max_bullets]
        # Add non-bullet content to first chunk only
        if i == 0:
            chunk_content = '\n'.join(non_bullet_lines + chunk_bullets)
        else:
            chunk_content = '\n'.join(chunk_bullets)
        bullet_chunks.append(chunk_content)
    
    return bullet_chunks

def parse_presentation_file(file_path):
    """Parse a markdown presentation file and extract slide data."""
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    
    # Extract title
    title_match = re.search(r'^# (.+?)$', content, re.MULTILINE)
    presentation_title = title_match.group(1) if title_match else "Presentation"
    
    # Extract slides
    slides = []
    slide_pattern = r'### (?:Slide|Diapositive) \d+: (.+?)\n\*\*(?:Layout|Mise en Page)\*\*: (.+?)\n\*\*(?:Content|Contenu)\*\*:\n(.*?)(?=\n\*Image Prompt|\n###|\Z)'
    
    for match in re.finditer(slide_pattern, content, re.DOTALL):
        slide_title = match.group(1)
        layout = match.group(2)
        content_text = match.group(3).strip()
        
        # Extract notes section if exists
        notes_text = None
        notes_match = re.search(r'\*\*(?:Notes|Remarques)\*\*:\n(.*?)(?=\n\*Image Prompt|\n###|\Z)', content[match.end():], re.DOTALL)
        if notes_match:
            notes_text = notes_match.group(1).strip()
        
        # Split content if it has too many bullets
        content_chunks = split_bullet_content(content_text)
        
        for i, chunk in enumerate(content_chunks):
            slide_data = {
                'title': slide_title if i == 0 else f"{slide_title} (Part {i+1})",
                'layout': layout,
                'content': chunk,
                'notes': notes_text if i == 0 else None  # Only add notes to first slide
            }
            slides.append(slide_data)
    
    return presentation_title, slides

def create_presentation_from_file(input_file_path, output_dir="findings_dir/presentations"):
    """Create a PowerPoint presentation from a parsed markdown file."""
    
    # Parse the input file
    presentation_title, slides = parse_presentation_file(input_file_path)
    
    # Create a presentation object
    prs = Presentation()
    
    for i, slide_data in enumerate(slides):
        # Choose layout based on slide type
        if slide_data['layout'] in ['Title Slide', 'Diapositive de Titre']:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            
            title.text = slide_data['title']
            if subtitle and slide_data['content']:
                # Extract subtitle from content
                lines = slide_data['content'].split('\n')
                subtitle_text = next((line.replace('- Sous-titre: ', '').replace('- Subtitle: ', '') for line in lines if 'Sous-titre:' in line or 'Subtitle:' in line), '')
                subtitle.text = subtitle_text
                
        else:
            # Content slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_data['title']
            content.text = slide_data['content']
        
        # Always add slide notes with content details
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        
        # Create detailed notes about the slide content
        notes_content = f"Slide: {slide_data['title']}\n\n"
        notes_content += f"Content Summary:\n{slide_data['content']}\n\n"
        
        # Add custom notes if they exist
        if slide_data.get('notes'):
            notes_content += f"Additional Notes:\n{slide_data['notes']}"
        
        text_frame.text = notes_content
    
    # Generate output filename
    safe_title = re.sub(r'[^\w\s-]', '', presentation_title).strip()
    safe_title = re.sub(r'[-\s]+', '_', safe_title)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    output_filename = f"{safe_title}_{timestamp}.pptx"
    
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, output_filename)
    
    # Save the presentation
    prs.save(output_path)
    print(f"Presentation saved as: {output_path}")
    return output_path

def main():
    """Main function to handle command line usage."""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python generate_dynamic_pptx.py <input_markdown_file> [output_directory]")
        print("Example: python generate_dynamic_pptx.py findings_dir/my_presentation.md")
        return
    
    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "findings_dir/presentations"
    
    if not os.path.exists(input_file):
        print(f"Error: Input file '{input_file}' not found.")
        return
    
    try:
        output_path = create_presentation_from_file(input_file, output_dir)
        print(f"Successfully created presentation: {output_path}")
    except Exception as e:
        print(f"Error creating presentation: {str(e)}")

if __name__ == "__main__":
    main()